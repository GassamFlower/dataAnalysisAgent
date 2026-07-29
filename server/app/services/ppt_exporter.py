"""PPT报告导出服务

职责：
- 将分析结果导出为PPT格式
- 生成7页结构化幻灯片
- 强制添加SIMULATED水印和免责声明

设计依据：docs/v-功能-PPT报告导出.md
"""
import io
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime


# 配色方案（与前端design tokens对齐）
COLOR_PRIMARY = RGBColor(139, 115, 85)  # #8B7355
COLOR_ACCENT = RGBColor(212, 165, 116)  # #D4A574
COLOR_BACKGROUND = RGBColor(250, 250, 248)  # #FAFAF8
COLOR_FOREGROUND = RGBColor(42, 42, 42)  # #2A2A2A
COLOR_GRAY = RGBColor(128, 128, 128)


def export_ppt(report_data: Dict[str, Any]) -> bytes:
    """导出PPT报告

    Args:
        report_data: 报告数据，包含信效度、诊断、差异检验等结果

    Returns:
        PPT文件的二进制内容
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    # 1. 封面页
    _add_cover_slide(prs, report_data)

    # 2. 信效度总览页
    _add_reliability_summary_slide(prs, report_data)

    # 3. 各维度信效度详情页
    _add_reliability_detail_slide(prs, report_data)

    # 4. 相关矩阵页
    _add_correlation_slide(prs, report_data)

    # 5. 差异检验页
    _add_diff_test_slide(prs, report_data)

    # 6. R4诊断页
    _add_diagnosis_slide(prs, report_data)

    # 7. 免责声明页
    _add_disclaimer_slide(prs)

    # 保存到字节流
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def _add_watermark(slide):
    """在幻灯片顶部添加SIMULATED水印"""
    left = Inches(0.5)
    top = Inches(0.25)
    width = Inches(12.333)
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "SIMULATED DATA - 仅供演示"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(192, 192, 192)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def _add_cover_slide(prs: Presentation, report_data: Dict[str, Any]):
    """添加封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11.333)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "数据分析报告"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    top = Inches(4.2)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    project_id = report_data.get("project_id", "N/A")
    p.text = f"项目ID: {project_id}"
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_FOREGROUND
    p.alignment = PP_ALIGN.CENTER

    # 生成时间
    top = Inches(5.2)
    height = Inches(0.6)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    p.text = f"生成时间: {now}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_GRAY
    p.alignment = PP_ALIGN.CENTER


def _add_reliability_summary_slide(prs: Presentation, report_data: Dict[str, Any]):
    """添加信效度总览页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(0.75)
    top = Inches(0.75)
    width = Inches(11.833)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "信效度总览"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 总体α
    overall_alpha = report_data.get("overall_alpha", 0)
    top = Inches(2)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"总量表 Cronbach's α: {overall_alpha:.3f}"
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_FOREGROUND

    # 通过维度数
    passed_count = report_data.get("passed_count", 0)
    total_count = report_data.get("total_count", 0)
    top = Inches(3.2)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"通过维度数: {passed_count} / {total_count}"
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_FOREGROUND


def _add_reliability_detail_slide(prs: Presentation, report_data: Dict[str, Any]):
    """添加各维度信效度详情页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(0.75)
    top = Inches(0.75)
    width = Inches(11.833)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "各维度信效度详情"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 表格
    reliability_results = report_data.get("reliability_results", [])
    if not reliability_results:
        return

    rows = len(reliability_results) + 1  # 表头 + 数据行
    cols = 5
    left = Inches(0.75)
    top = Inches(2)
    width = Inches(11.833)
    height = Inches(0.8) * rows

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 表头
    headers = ["维度", "Cronbach's α", "KMO", "Bartlett p", "是否通过"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ACCENT

    # 数据行
    for row_idx, result in enumerate(reliability_results, start=1):
        table.cell(row_idx, 0).text = result.get("dimension", "")
        table.cell(row_idx, 1).text = f"{result.get('alpha', 0):.3f}"
        table.cell(row_idx, 2).text = f"{result.get('kmo', 0):.3f}"
        table.cell(row_idx, 3).text = f"{result.get('bartlett_p_value', 0):.5f}"
        passed = "✓" if result.get("passed") else "✗"
        table.cell(row_idx, 4).text = passed

        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text_frame.paragraphs[0].font.size = Pt(12)


def _add_correlation_slide(prs: Presentation, report_data: Dict[str, Any]):
    """添加相关矩阵页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(0.75)
    top = Inches(0.75)
    width = Inches(11.833)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "相关分析"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 说明文本
    top = Inches(2)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "相关矩阵分析结果（如有）"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_FOREGROUND


def _add_diff_test_slide(prs: Presentation, report_data: Dict[str, Any]):
    """添加差异检验页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(0.75)
    top = Inches(0.75)
    width = Inches(11.833)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "差异检验"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 差异检验结果
    diff_tests = report_data.get("diff_tests", [])
    if not diff_tests:
        top = Inches(2)
        height = Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "未配置假设路径，无差异检验结果"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_FOREGROUND
        return

    # 表格
    rows = len(diff_tests) + 1
    cols = 5
    left = Inches(0.75)
    top = Inches(2)
    width = Inches(11.833)
    height = Inches(0.8) * rows

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 表头
    headers = ["假设路径", "检验方法", "统计量", "p值", "显著性"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ACCENT

    # 数据行
    for row_idx, test in enumerate(diff_tests, start=1):
        predictor = test.get("predictor", "")
        outcome = test.get("outcome", "")
        path = f"{predictor} → {outcome}"
        table.cell(row_idx, 0).text = path
        table.cell(row_idx, 1).text = test.get("method_name", "")
        table.cell(row_idx, 2).text = f"{test.get('statistic', 0):.3f}"
        table.cell(row_idx, 3).text = f"{test.get('p_value', 0):.5f}"
        significant = "显著 *" if test.get("significant") else "不显著"
        table.cell(row_idx, 4).text = significant

        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text_frame.paragraphs[0].font.size = Pt(12)


def _add_diagnosis_slide(prs: Presentation, report_data: Dict[str, Any]):
    """添加R4诊断页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(0.75)
    top = Inches(0.75)
    width = Inches(11.833)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "R4 诊断结论"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 诊断结果
    diagnosis = report_data.get("diagnosis")
    if not diagnosis:
        top = Inches(2)
        height = Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "无诊断数据"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_FOREGROUND
        return

    passed = diagnosis.get("passed", False)
    status_text = "通过" if passed else "不通过"

    top = Inches(2)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"诊断结果: {status_text}"
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_FOREGROUND

    # 问题列表
    issues = diagnosis.get("issues", [])
    if issues:
        top = Inches(3)
        height = Inches(0.6)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"发现问题数: {len(issues)}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_FOREGROUND

        # 问题详情（最多显示3个）
        top = Inches(3.8)
        for idx, issue in enumerate(issues[:3], start=1):
            height = Inches(0.8)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            dimension = issue.get("dimension", "")
            metric = issue.get("metric", "")
            reason = issue.get("reason", "")
            p.text = f"{idx}. {dimension} - {metric}: {reason}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_FOREGROUND
            top += Inches(0.8)


def _add_disclaimer_slide(prs: Presentation):
    """添加免责声明页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 水印
    _add_watermark(slide)

    # 标题
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11.333)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "免责声明"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # 免责声明内容
    top = Inches(4)
    height = Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = (
        "本报告基于模拟数据生成，仅供学术研究和教学演示使用，"
        "不代表真实统计分析结果。实际数据分析请联系专业统计人员。"
    )
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_GRAY
    p.alignment = PP_ALIGN.CENTER
    text_frame.word_wrap = True
