"""PPT 报告导出服务单元测试。

覆盖（testing-strategy 第二步：三类用例）：
- 正常：完整数据导出、9 页结构验证、水印验证
- 边界：空数据、缺失字段、空列表
- 异常：缺少 python-pptx 依赖（importorskip）

设计依据：docs/v-功能-PPT报告导出.md
"""
import io

import pytest

# python-pptx 为可选依赖，未安装时跳过全部测试
pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

from app.services.ppt_exporter import export_ppt  # noqa: E402


def _extract_all_text(slide) -> str:
    """提取幻灯片中所有文本（含文本框和表格单元格）。"""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return " ".join(parts)


def _build_full_report_data():
    """构造完整的报告数据（含信效度、差异检验、诊断）。"""
    return {
        "project_id": "test-project-001",
        "overall_alpha": 0.856,
        "passed_count": 4,
        "total_count": 5,
        "reliability_results": [
            {
                "dimension": "学习动机",
                "alpha": 0.82,
                "kmo": 0.75,
                "bartlett_p_value": 0.001,
                "passed": True,
            },
            {
                "dimension": "自我效能",
                "alpha": 0.65,
                "kmo": 0.70,
                "bartlett_p_value": 0.005,
                "passed": False,
            },
        ],
        "diff_tests": [
            {
                "predictor": "学习动机",
                "outcome": "学业成绩",
                "method_name": "Pearson 相关",
                "statistic": 0.456,
                "p_value": 0.00012,
                "significant": True,
            },
        ],
        "diagnosis": {
            "passed": False,
            "issues": [
                {"dimension": "自我效能", "metric": "Cronbach's α", "reason": "α < 0.7"},
            ],
        },
        "sample_representativeness": {
            "grade": "B",
            "overall_score": 82,
            "sample_size": 120,
            "summary": "样本结构良好，接近目标总体。",
            "items": [
                {
                    "title": "样本量（N）",
                    "status": "pass",
                    "message": "样本量充足。",
                },
                {
                    "title": "性别分布",
                    "status": "warn",
                    "message": "性别比例略有偏差。",
                },
            ],
        },
        "sample_size_plan": {
            "analysis_type": "correlation",
            "effect_size": 0.3,
            "effect_source": "default",
            "required_n": 85,
            "per_group_n": None,
            "recommended_n": 100,
            "planned_n": 120,
            "verdict": "sufficient",
            "verdict_label": "已达标",
            "guidance": ["样本量已达标，可放心进行统计分析。"],
        },
    }


# ============================
# 正常用例
# ============================

def test_export_ppt_returns_valid_bytes():
    """正常：导出返回有效的 PPT 二进制内容。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    assert isinstance(result, bytes)
    assert len(result) > 0
    # PPT 文件以 PK 开头（ZIP 容器）
    assert result[:2] == b"PK"


def test_export_ppt_has_nine_slides():
    """正常：导出的 PPT 包含 9 页幻灯片（封面/信效度总览/详情/相关/差异/诊断/代表性/规划/免责声明）。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 9


def test_export_ppt_watermark_on_every_slide():
    """正常：每页幻灯片都包含 SIMULATED 水印。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    for idx, slide in enumerate(prs.slides):
        all_text = _extract_all_text(slide)
        assert "SIMULATED" in all_text, f"第 {idx + 1} 页缺少 SIMULATED 水印"


def test_export_ppt_cover_slide_content():
    """正常：封面页包含标题和项目 ID。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    cover_slide = prs.slides[0]
    cover_text = _extract_all_text(cover_slide)
    assert "数据分析报告" in cover_text
    assert "test-project-001" in cover_text


def test_export_ppt_reliability_summary_values():
    """正常：信效度总览页正确显示 α 值和通过维度数。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    summary_slide = prs.slides[1]
    summary_text = _extract_all_text(summary_slide)
    assert "0.856" in summary_text  # overall_alpha
    assert "4" in summary_text and "5" in summary_text  # passed_count / total_count


def test_export_ppt_reliability_detail_table():
    """正常：信效度详情页表格包含所有维度数据。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    detail_slide = prs.slides[2]
    detail_text = _extract_all_text(detail_slide)
    assert "学习动机" in detail_text
    assert "自我效能" in detail_text
    assert "0.820" in detail_text  # alpha for 学习动机


def test_export_ppt_diff_test_table():
    """正常：差异检验页表格包含检验结果。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    diff_slide = prs.slides[4]
    diff_text = _extract_all_text(diff_slide)
    assert "学习动机" in diff_text
    assert "学业成绩" in diff_text
    assert "Pearson" in diff_text


def test_export_ppt_diagnosis_slide():
    """正常：诊断页显示诊断结果和问题列表。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    diag_slide = prs.slides[5]
    diag_text = _extract_all_text(diag_slide)
    assert "不通过" in diag_text  # diagnosis passed=False
    assert "自我效能" in diag_text  # issue dimension


def test_export_ppt_disclaimer_slide():
    """正常：免责声明页包含免责声明文本。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    disclaimer_slide = prs.slides[8]
    disclaimer_text = _extract_all_text(disclaimer_slide)
    assert "免责声明" in disclaimer_text
    assert "模拟数据" in disclaimer_text


def test_export_ppt_representativeness_slide():
    """正常：样本代表性页显示评级与检查项。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    rep_slide = prs.slides[6]
    rep_text = _extract_all_text(rep_slide)
    assert "样本代表性诊断" in rep_text
    assert "B" in rep_text and "82" in rep_text  # grade / score
    assert "样本量充足" in rep_text  # item message


def test_export_ppt_sample_plan_slide():
    """正常：样本量规划页显示目标、已收 N 与达成判定。"""
    data = _build_full_report_data()
    result = export_ppt(data)

    prs = Presentation(io.BytesIO(result))
    plan_slide = prs.slides[7]
    plan_text = _extract_all_text(plan_slide)
    assert "样本量规划与回收目标" in plan_text
    assert "100" in plan_text  # recommended_n
    assert "120" in plan_text  # planned_n（已收 N）
    assert "已达标" in plan_text  # verdict_label


# ============================
# 边界用例
# ============================

def test_export_ppt_empty_data():
    """边界：空数据字典仍能导出 9 页（使用默认值）。"""
    result = export_ppt({})

    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 9

    # 封面页项目 ID 应为 N/A
    cover_text = _extract_all_text(prs.slides[0])
    assert "N/A" in cover_text


def test_export_ppt_empty_reliability_results():
    """边界：信效度结果为空列表时，详情页不应崩溃。"""
    data = _build_full_report_data()
    data["reliability_results"] = []

    result = export_ppt(data)
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 9


def test_export_ppt_empty_diff_tests():
    """边界：差异检验为空列表时，差异页显示提示文本。"""
    data = _build_full_report_data()
    data["diff_tests"] = []

    result = export_ppt(data)
    prs = Presentation(io.BytesIO(result))

    diff_slide = prs.slides[4]
    diff_text = _extract_all_text(diff_slide)
    assert "未配置假设路径" in diff_text


def test_export_ppt_no_diagnosis():
    """边界：诊断数据为 None 时，诊断页显示'无诊断数据'。"""
    data = _build_full_report_data()
    data["diagnosis"] = None

    result = export_ppt(data)
    prs = Presentation(io.BytesIO(result))

    diag_slide = prs.slides[5]
    diag_text = _extract_all_text(diag_slide)
    assert "无诊断数据" in diag_text


def test_export_ppt_no_representativeness():
    """边界：无代表性数据时，代表性页显示提示文本。"""
    data = _build_full_report_data()
    data["sample_representativeness"] = None

    result = export_ppt(data)
    prs = Presentation(io.BytesIO(result))

    rep_slide = prs.slides[6]
    rep_text = _extract_all_text(rep_slide)
    assert "无样本代表性数据" in rep_text


def test_export_ppt_no_sample_plan():
    """边界：无规划数据时，规划页显示提示文本。"""
    data = _build_full_report_data()
    data["sample_size_plan"] = None

    result = export_ppt(data)
    prs = Presentation(io.BytesIO(result))

    plan_slide = prs.slides[7]
    plan_text = _extract_all_text(plan_slide)
    assert "无样本量规划数据" in plan_text
